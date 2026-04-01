"""Generate Word + PowerPoint proposal documents for Space ONE.Z & Old Beat Yeongdo website project."""

from docx import Document
from docx.shared import Pt, RGBColor, Cm, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN
import os

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))


# ══════════════════════════════════════════════════════════════
#  WORD — 웹사이트 제작 제안서
# ══════════════════════════════════════════════════════════════
def generate_word():
    doc = Document()

    style = doc.styles["Normal"]
    style.font.size = Pt(11)
    style.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
    style.paragraph_format.line_spacing = 1.6

    def heading(text, level=1):
        h = doc.add_heading(text, level=level)
        for r in h.runs:
            r.font.color.rgb = RGBColor(0x05, 0x05, 0x05)

    def centered(text, size=11, bold=False, color=None, italic=False):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = p.add_run(text)
        r.font.size = Pt(size)
        r.bold = bold
        r.italic = italic
        if color:
            r.font.color.rgb = RGBColor(*color)

    def add_table(headers, rows):
        t = doc.add_table(rows=len(rows) + 1, cols=len(headers))
        t.style = "Light Grid Accent 1"
        for i, h in enumerate(headers):
            cell = t.rows[0].cells[i]
            cell.text = h
            for p in cell.paragraphs:
                for r in p.runs:
                    r.bold = True
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                t.rows[ri + 1].cells[ci].text = val
        doc.add_paragraph()

    # ═══════════════════════════════════════════
    # 표지
    # ═══════════════════════════════════════════
    for _ in range(5):
        doc.add_paragraph()
    centered("웹사이트 제작 제안서", 32, True, (0xC8, 0x96, 0x5A))
    doc.add_paragraph()
    centered("스페이스원지 & 올드비트 영도", 22, True, (0x1A, 0x1A, 0x1A))
    centered("Space ONE.Z & Old Beat Yeongdo", 14, False, (0x99, 0x99, 0x99), True)
    doc.add_paragraph()
    centered("공식 웹사이트 구축 프로젝트", 16, False, (0x55, 0x55, 0x55))

    for _ in range(6):
        doc.add_paragraph()

    # 제안 정보 표
    info_table = doc.add_table(rows=5, cols=2)
    info_table.alignment = WD_TABLE_ALIGNMENT.CENTER
    info_data = [
        ("제안사", "Skyline Works (스카이라인웍스)"),
        ("제안 대상", "스페이스원지 & 올드비트 영도"),
        ("제안일", "2025년 4월"),
        ("문서 버전", "v1.0"),
        ("담당자", "Skyline Works 프로젝트팀"),
    ]
    for i, (label, value) in enumerate(info_data):
        info_table.rows[i].cells[0].text = label
        info_table.rows[i].cells[1].text = value
        for p in info_table.rows[i].cells[0].paragraphs:
            p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            for r in p.runs:
                r.font.size = Pt(10)
                r.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        for p in info_table.rows[i].cells[1].paragraphs:
            for r in p.runs:
                r.font.size = Pt(10)
                r.bold = True

    doc.add_page_break()

    # ═══════════════════════════════════════════
    # 목차
    # ═══════════════════════════════════════════
    heading("목차")
    toc_items = [
        "1. 제안 개요",
        "2. 제안사 소개",
        "3. 프로젝트 이해",
        "4. 웹사이트 설계 방안",
        "5. 디자인 컨셉",
        "6. 기술 구현 방안",
        "7. 세부 기능 명세",
        "8. 추진 일정",
        "9. 투입 인력 및 조직",
        "10. 견적 및 비용",
        "11. 유지보수 및 운영 지원",
        "12. 기대 효과",
        "13. 별첨: 포트폴리오",
    ]
    for item in toc_items:
        p = doc.add_paragraph(item)
        p.paragraph_format.space_after = Pt(4)
    doc.add_page_break()

    # ═══════════════════════════════════════════
    # 1. 제안 개요
    # ═══════════════════════════════════════════
    heading("1. 제안 개요")
    add_table(["항목", "내용"], [
        ["프로젝트명", "스페이스원지 & 올드비트 영도 공식 웹사이트 구축"],
        ["프로젝트 유형", "신규 웹사이트 기획 · 디자인 · 개발"],
        ["클라이언트", "스페이스원지 & 올드비트 영도"],
        ["위치", "부산광역시 영도구 봉래나루로 214 1층"],
        ["공간 유형", "복합문화공간 (레스토랑 + 갤러리 + 공연장 + 워크숍)"],
        ["웹사이트 형태", "반응형 원페이지 + 서브페이지 (총 8개 페이지)"],
        ["예약 연동", "네이버 예약 외부 링크"],
        ["목표 런칭일", "착수 후 5주"],
    ])

    doc.add_paragraph(
        "본 제안서는 부산 영도의 복합문화공간 '스페이스원지 & 올드비트 영도'의 "
        "브랜드 가치를 온라인에 효과적으로 전달하고, 방문 예약과 대관 문의를 "
        "원활히 수행할 수 있는 공식 웹사이트 구축 프로젝트에 대한 제안입니다."
    )
    doc.add_page_break()

    # ═══════════════════════════════════════════
    # 2. 제안사 소개
    # ═══════════════════════════════════════════
    heading("2. 제안사 소개")
    heading("2.1 회사 개요", 2)
    add_table(["항목", "내용"], [
        ["회사명", "Skyline Works (스카이라인웍스)"],
        ["사업 영역", "비즈니스 로직 설계 · 웹/앱 플랫폼 개발 · AI 솔루션"],
        ["기술 스택", "Next.js, React, TypeScript, Tailwind CSS, Supabase, AI"],
        ["핵심 슬로건", "예쁜 웹사이트는 필요 없습니다. 우리는 팔리는 논리를 만듭니다."],
    ])

    heading("2.2 핵심 역량", 2)
    for item in [
        "Solid Logic Structure — 단순한 기능 나열이 아닌, 고객이 왜 이 서비스를 써야 하는지 설득하는 비즈니스 시나리오 설계",
        "Irresistible Flow — 고객의 시선 이동과 심리를 계산한 자연스러운 전환 UX 설계",
        "Visual Persuasion — 설계된 논리를 시각적으로 각인시키는 압도적인 디자인 구현",
        "Communication Bridge — 비즈니스 언어를 기술 언어로 완벽히 통역",
        "Operational Intelligence — 런칭 이후 운영까지 고려한 지속 가능한 구조 설계",
    ]:
        doc.add_paragraph(item, style="List Bullet")

    heading("2.3 차별화 포인트", 2)
    add_table(["기존 개발사", "Skyline Works"], [
        ["'기능이 잘 작동하는가?'", "'이 구성이 고객을 설득할 수 있는가?'"],
        ["'디자인이 예쁜가?'", "'이 논리가 비즈니스를 지탱하는가?'"],
        ["기능·심미성 중심", "설득·논리 중심"],
        ["개발만 하고 끝", "런칭 이후 운영·정산·행정까지 고려"],
    ])
    doc.add_page_break()

    # ═══════════════════════════════════════════
    # 3. 프로젝트 이해
    # ═══════════════════════════════════════════
    heading("3. 프로젝트 이해")
    heading("3.1 사업 배경", 2)
    doc.add_paragraph(
        "부산 영도는 대한민국 근현대 조선산업의 중심지로, 최근 도시재생 사업과 함께 "
        "문화예술 공간으로 변모하고 있습니다. 스페이스원지와 올드비트 영도는 영도의 "
        "역사적 산업 건물을 리노베이션하여 레스토랑, 갤러리, 공연장, 워크숍 공간을 "
        "갖춘 복합문화공간으로 탄생했습니다."
    )
    doc.add_paragraph(
        "이 공간의 가치와 프로그램을 효과적으로 전달하고, 온라인을 통한 예약과 "
        "대관 문의를 가능하게 하는 공식 웹사이트 구축이 필요한 시점입니다."
    )

    heading("3.2 핵심 과제", 2)
    for item in [
        "브랜드 아이덴티티 확립: 두 브랜드(스페이스원지 + 올드비트)의 독창적 이미지를 하나의 웹사이트에 통합 구현",
        "정보 허브 구축: 공간 소개, 메뉴, 행사/전시 일정, 소식 등 모든 정보를 한곳에서 제공",
        "예약 전환 최적화: 네이버 예약 연동을 통한 원활한 레스토랑/행사 예약 유도",
        "대관 문의 채널 확보: 공간 대관 희망자를 위한 체계적 온라인 문의 시스템",
        "문화 콘텐츠 플랫폼: 전시, 공연, 워크숍 등 문화 프로그램의 지속적 홍보 채널 역할",
    ]:
        doc.add_paragraph(item, style="List Bullet")

    heading("3.3 주요 타겟", 2)
    add_table(["타겟 그룹", "설명", "핵심 니즈"], [
        ["미식 애호가", "영도 로컬 식재료 기반 코스 다이닝에 관심 있는 2030~4050세대", "메뉴/예약 정보"],
        ["문화예술 관심층", "전시, 공연, 워크숍에 관심 있는 예술 애호가 및 크리에이터", "행사 일정/참여 방법"],
        ["관광객", "부산 여행 중 영도의 문화 명소를 찾는 국내외 관광객", "접근성/공간 안내"],
        ["기업/단체", "기업 행사, 파티, 촬영 등 공간 대관을 원하는 기업 및 단체", "대관 정보/문의"],
        ["지역 커뮤니티", "영도 및 부산 지역 주민, 로컬 아티스트", "프로그램 참여"],
    ])
    doc.add_page_break()

    # ═══════════════════════════════════════════
    # 4. 웹사이트 설계 방안
    # ═══════════════════════════════════════════
    heading("4. 웹사이트 설계 방안")
    heading("4.1 사이트맵", 2)
    add_table(["경로", "페이지명", "주요 콘텐츠", "비고"], [
        ["/", "홈", "히어로 + 공간소개 + 인터루드 + 행사하이라이트 + 최신소식", "랜딩 페이지"],
        ["/about", "소개", "브랜드 스토리, 운영 철학, 통계", "브랜드 아이덴티티"],
        ["/news", "소식", "소식/공지/미디어 카드 그리드", "콘텐츠 허브"],
        ["/news/[slug]", "소식 상세", "개별 소식 본문", "동적 라우팅"],
        ["/events", "행사/전시", "진행중/예정/종료 분류 목록", "상태별 필터"],
        ["/events/[slug]", "행사 상세", "행사 본문 + 예약 CTA", "동적 라우팅"],
        ["/space", "공간 안내", "레스토랑, 갤러리, 루프탑, 워크숍룸", "4개 공간"],
        ["/reservation", "예약/대관", "네이버 예약 링크 + 대관 문의 폼", "전환 페이지"],
    ])

    heading("4.2 페이지별 상세 설계", 2)
    pages_detail = {
        "홈페이지 (/)": [
            "풀스크린 히어로 (120vh): 패럴랙스 배경 이미지 + 줌인 애니메이션 + 대형 타이포그래피",
            "공간 소개 (SpaceIntro): 풀스크린 시네마틱 섹션 x2, 패럴랙스 배경 + 좌/우 슬라이드인 텍스트",
            "인터루드 (Interlude): 풀블리드 인용문, 스케일업 리빌 - '오래된 것의 새로운 울림'",
            "행사 하이라이트 (EventHighlight): 21:9 피처드 카드 + 이벤트 그리드",
            "최신 소식 (LatestNews): 에디토리얼 리스트 (번호 + 호버 밀림 효과)",
        ],
        "소개 (/about)": [
            "히어로: 21:9 이미지 + 카피라이팅",
            "Our Story: 영도 + 조선업 역사 + 리노베이션 스토리",
            "Philosophy: 운영 철학 - '맛과 예술, 그리고 사람이 만나는 곳'",
            "통계 섹션: 오픈연도(2025), 문화공간(4개), 365일의 이야기",
            "CTA: 공간 안내 + 예약하기 버튼",
        ],
        "소식 (/news)": [
            "3열 카드 그리드 레이아웃",
            "카테고리: 공지 / 소식 / 미디어",
            "썸네일 + 제목 + 요약 + 날짜",
            "상세 페이지: slug 기반 동적 라우팅",
        ],
        "행사/전시 (/events)": [
            "상태별 분류: 진행중 / 예정 / 종료",
            "카테고리: 전시 / 공연 / 워크숍 / 팝업",
            "상세 페이지: 행사 본문 + 장소 + 일정",
            "네이버 예약 CTA 버튼",
        ],
        "공간 안내 (/space)": [
            "4개 공간 소개: 레스토랑(40석), 갤러리(80명), 루프탑(60명), 워크숍룸(20명)",
            "이미지 그리드 + 특징 태그",
            "좌우 교차 레이아웃 (지그재그)",
            "대관 문의 CTA",
        ],
        "예약/대관 (/reservation)": [
            "레스토랑 예약: 네이버 예약 외부 링크 + 런치/디너 코스 메뉴 카드",
            "행사 예약: 네이버 예약 연동",
            "대관 문의 폼: 이름, 연락처, 이메일, 희망일자, 행사유형(7종), 예상인원, 상세내용",
            "접수 확인 UI: 체크마크 + 안내 메시지",
        ],
    }
    for title, items in pages_detail.items():
        p = doc.add_paragraph()
        r = p.add_run(title)
        r.bold = True
        r.font.size = Pt(12)
        for item in items:
            doc.add_paragraph(item, style="List Bullet")
    doc.add_page_break()

    # ═══════════════════════════════════════════
    # 5. 디자인 컨셉
    # ═══════════════════════════════════════════
    heading("5. 디자인 컨셉")
    heading("5.1 디자인 방향", 2)
    doc.add_paragraph(
        "\"Warm Industrial\" — 풀블리드 이미지, 패럴랙스 스크롤, 대형 타이포그래피, "
        "GPU 가속 애니메이션을 활용한 몰입감 있는 시네마틱 디자인. "
        "영도 조선소 유산의 따뜻한 질감과 앤틱 브래스 톤으로 문화공간·다이닝의 감성을 표현합니다."
    )

    heading("5.2 컬러 팔레트", 2)
    add_table(["요소", "색상 코드", "설명", "사용 위치"], [
        ["배경", "#0A0A08", "웜 블랙", "전체 배경"],
        ["액센트", "#C8965A", "앤틱 브래스", "CTA, 강조, 태그"],
        ["메인 텍스트", "#F0EBE3", "웜 크림", "제목, 본문"],
        ["서브 텍스트", "#8A8578", "웜 그레이", "설명, 보조 텍스트"],
        ["표면", "#171612", "다크 웜", "카드, 입력 배경"],
    ])

    heading("5.3 타이포그래피", 2)
    doc.add_paragraph("Geist Sans (영문) + Pretendard (한글) — 모던하고 기하학적인 서체 조합")

    heading("5.4 인터랙션 & 애니메이션", 2)
    add_table(["이름", "동작", "사용 위치"], [
        ["fade-up", "아래에서 위로 페이드인 (60px)", "기본 섹션 진입"],
        ["reveal-scale", "0.92 → 1 스케일업 페이드인", "인터루드, 피처드 콘텐츠"],
        ["reveal-left/right", "좌/우 80px 슬라이드인", "공간 소개 교차 레이아웃"],
        ["패럴랙스", "스크롤 기반 translateY (0.4배속)", "히어로, 공간 배경 이미지"],
        ["hero-zoom", "1.15 → 1 스케일 (8초)", "히어로 배경 진입"],
        ["bounce-slow", "스크롤 인디케이터 바운스", "히어로 하단"],
    ])
    doc.add_page_break()

    # ═══════════════════════════════════════════
    # 6. 기술 구현 방안
    # ═══════════════════════════════════════════
    heading("6. 기술 구현 방안")
    heading("6.1 기술 스택", 2)
    add_table(["영역", "기술", "선정 사유"], [
        ["프레임워크", "Next.js 16 (App Router)", "SSG 정적 생성 + SEO 최적화 + 최신 React 지원"],
        ["UI 라이브러리", "React 19 + TypeScript", "Server Components + 타입 안전성"],
        ["스타일링", "Tailwind CSS v4", "유틸리티 CSS + 커스텀 테마 + 빠른 개발"],
        ["애니메이션", "CSS Transitions + IntersectionObserver", "네이티브 성능 + GPU 가속"],
        ["스크롤", "Lenis (Smooth Scroll)", "부드러운 스크롤 경험"],
        ["콘텐츠 관리", "JSON 파일 기반", "Git push 시 자동 반영 + CMS 불필요"],
        ["예약 시스템", "네이버 예약 외부 링크", "기존 플랫폼 활용 + 추가 개발 불필요"],
        ["배포", "Vercel", "자동 빌드/배포 + CDN + SSL"],
    ])

    heading("6.2 콘텐츠 관리 구조", 2)
    doc.add_paragraph(
        "별도의 CMS 없이 JSON 파일 기반으로 콘텐츠를 관리합니다. "
        "Git에 push하면 Vercel이 자동으로 빌드/배포하므로, "
        "콘텐츠 업데이트가 간편하고 추가 비용이 발생하지 않습니다."
    )
    add_table(["데이터 파일", "관리 콘텐츠", "주요 필드"], [
        ["news.json", "소식/공지/미디어", "제목, 날짜, 썸네일, 요약, 본문, 카테고리"],
        ["events.json", "행사/전시 프로그램", "제목, 시작/종료일, 상태, 장소, 카테고리"],
        ["spaces.json", "공간 정보", "공간명, 설명, 이미지, 수용인원, 특징"],
    ])

    heading("6.3 성능 최적화", 2)
    for item in [
        "SSG (Static Site Generation): 빌드 시 HTML 사전 생성으로 최고 속도 달성",
        "이미지 최적화: Next.js Image 컴포넌트 활용 (WebP, lazy loading)",
        "GPU 가속 애니메이션: transform3d + will-change로 60fps 보장",
        "코드 스플리팅: 페이지별 자동 번들 분리로 초기 로딩 최소화",
    ]:
        doc.add_paragraph(item, style="List Bullet")
    doc.add_page_break()

    # ═══════════════════════════════════════════
    # 7. 세부 기능 명세
    # ═══════════════════════════════════════════
    heading("7. 세부 기능 명세")

    heading("7.1 공통 기능", 2)
    add_table(["기능", "상세 내용"], [
        ["반응형 디자인", "320px~1920px 완전 반응형, 모바일/태블릿/데스크탑 최적화"],
        ["모바일 네비게이션", "햄버거 메뉴 + 슬라이드 패널"],
        ["SEO 최적화", "Next.js Metadata API, 페이지별 title/description, OG 태그"],
        ["스무스 스크롤", "Lenis 라이브러리 기반 부드러운 스크롤"],
        ["스크롤 애니메이션", "IntersectionObserver 기반 요소별 리빌 애니메이션"],
        ["Header/Footer", "투명 Header + 스크롤 시 배경 전환, 연락처/SNS Footer"],
    ])

    heading("7.2 네이버 예약 연동", 2)
    for item in [
        "레스토랑: 런치 코스(5만원) / 디너 코스(8만원) 메뉴 카드 + CTA 버튼",
        "행사: 재즈 나이트, 워크숍 등 행사별 예약 버튼",
        "네이버 예약 페이지로 외부 링크 이동 (새 탭)",
    ]:
        doc.add_paragraph(item, style="List Bullet")

    heading("7.3 대관 문의 폼", 2)
    add_table(["필드", "타입", "필수"], [
        ["이름", "텍스트", "O"],
        ["연락처", "전화번호", "O"],
        ["이메일", "이메일", "O"],
        ["희망 일자", "날짜 선택", "O"],
        ["행사 유형", "드롭다운 (7종)", "O"],
        ["예상 인원", "숫자", ""],
        ["상세 내용", "장문 텍스트", ""],
    ])
    doc.add_paragraph("행사 유형: 전시, 공연/콘서트, 워크숍/클래스, 기업 행사, 프라이빗 파티, 촬영, 기타")
    doc.add_page_break()

    # ═══════════════════════════════════════════
    # 8. 추진 일정
    # ═══════════════════════════════════════════
    heading("8. 추진 일정")
    doc.add_paragraph("총 5주 과정으로 진행하며, 각 단계별 산출물을 제출합니다.")
    add_table(["주차", "단계", "주요 작업", "산출물"], [
        ["1주차", "프로젝트 셋업", "개발 환경 구성, 디자인 토큰 정의, Header/Footer 구현, 공통 컴포넌트", "프로젝트 기본 구조"],
        ["2주차", "홈페이지 구현", "히어로, 공간소개, 인터루드, 행사하이라이트, 최신소식 섹션", "홈페이지 시안 리뷰"],
        ["3주차", "서브 페이지 구현", "소개, 소식(목록+상세), 행사(목록+상세), 공간안내, 예약/대관", "전체 페이지 시안 리뷰"],
        ["4주차", "콘텐츠 + QA", "실제 사진/텍스트 적용, 크로스브라우저 테스트, 성능 최적화", "테스트 완료 보고"],
        ["5주차", "배포 + 런칭", "Vercel 배포, 커스텀 도메인, 최종 점검, 인수인계", "최종 산출물 인도"],
    ])
    doc.add_page_break()

    # ═══════════════════════════════════════════
    # 9. 투입 인력 및 조직
    # ═══════════════════════════════════════════
    heading("9. 투입 인력 및 조직")
    add_table(["역할", "인원", "주요 업무", "투입 기간"], [
        ["프로젝트 매니저 (PM)", "1명", "일정 관리, 고객 커뮤니케이션, 품질 관리", "전 기간"],
        ["UI/UX 디자이너", "1명", "디자인 컨셉, 시안 제작, 인터랙션 설계", "1~3주차"],
        ["프론트엔드 개발자", "1명", "Next.js 개발, 반응형 구현, 애니메이션", "전 기간"],
        ["QA 엔지니어", "1명", "크로스브라우저, 성능, 접근성 테스트", "4~5주차"],
    ])
    doc.add_page_break()

    # ═══════════════════════════════════════════
    # 10. 견적 및 비용
    # ═══════════════════════════════════════════
    heading("10. 견적 및 비용")
    heading("10.1 개발 비용", 2)
    add_table(["항목", "내용", "금액 (부가세 별도)"], [
        ["기획 및 설계", "사이트맵, 와이어프레임, 콘텐츠 구조 설계", "1,500,000원"],
        ["UI/UX 디자인", "디자인 컨셉, 시안, 인터랙션 설계", "2,500,000원"],
        ["프론트엔드 개발", "8개 페이지 개발, 반응형, 애니메이션", "4,000,000원"],
        ["콘텐츠 적용", "실제 이미지/텍스트 적용, 최적화", "500,000원"],
        ["테스트 및 배포", "QA, Vercel 배포, 도메인 연결", "500,000원"],
        ["", "", ""],
        ["합계", "", "9,000,000원"],
        ["부가세 (10%)", "", "900,000원"],
        ["총 금액", "", "9,900,000원"],
    ])

    heading("10.2 결제 조건", 2)
    add_table(["단계", "시점", "비율", "금액"], [
        ["계약금", "계약 체결 시", "40%", "3,960,000원"],
        ["중도금", "시안 확인 후 (3주차)", "30%", "2,970,000원"],
        ["잔금", "최종 납품 완료 후", "30%", "2,970,000원"],
    ])

    doc.add_paragraph()
    p = doc.add_paragraph()
    r = p.add_run("* 상기 금액은 제안 기준이며, 상세 협의를 통해 조정될 수 있습니다.")
    r.font.size = Pt(9)
    r.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    r.italic = True
    doc.add_page_break()

    # ═══════════════════════════════════════════
    # 11. 유지보수 및 운영 지원
    # ═══════════════════════════════════════════
    heading("11. 유지보수 및 운영 지원")

    heading("11.1 무상 유지보수", 2)
    for item in [
        "납품 후 3개월간 무상 하자보수 제공",
        "버그 수정, 브라우저 호환성 이슈 대응",
        "경미한 텍스트/이미지 수정 (월 2회 이내)",
    ]:
        doc.add_paragraph(item, style="List Bullet")

    heading("11.2 콘텐츠 업데이트 가이드", 2)
    for item in [
        "JSON 파일 수정 가이드 문서 제공",
        "소식/행사 업데이트 매뉴얼",
        "이미지 규격 및 최적화 가이드",
    ]:
        doc.add_paragraph(item, style="List Bullet")

    heading("11.3 유상 유지보수 (선택)", 2)
    add_table(["플랜", "내용", "월 비용"], [
        ["Basic", "월 4회 콘텐츠 업데이트 + 모니터링", "200,000원/월"],
        ["Standard", "월 8회 업데이트 + 기능 개선 + 분석 리포트", "400,000원/월"],
        ["Premium", "상시 대응 + 신규 기능 개발 + SEO 관리", "별도 협의"],
    ])
    doc.add_page_break()

    # ═══════════════════════════════════════════
    # 12. 기대 효과
    # ═══════════════════════════════════════════
    heading("12. 기대 효과")
    effects = [
        ("브랜드 인지도 향상", "시네마틱 웹사이트를 통해 영도의 대표 문화공간으로서의 브랜드 이미지를 온라인에 확립합니다."),
        ("방문자 유입 증가", "SEO 최적화와 SNS 공유를 통해 온라인 노출을 확대하고, 실제 방문으로 전환합니다."),
        ("예약 전환율 향상", "네이버 예약과의 원활한 연동으로 검색에서 예약까지 원스톱 경험을 제공합니다."),
        ("대관 문의 효율화", "온라인 문의 폼을 통해 체계적으로 접수하고 관리할 수 있습니다."),
        ("문화 프로그램 홍보", "전시/공연/워크숍 정보를 지속적으로 업데이트하여 재방문을 유도합니다."),
        ("운영 효율 극대화", "JSON 기반 콘텐츠 관리로 별도의 CMS 비용 없이 간편하게 업데이트할 수 있습니다."),
    ]
    for title, desc in effects:
        p = doc.add_paragraph()
        r = p.add_run(f"{title}: ")
        r.bold = True
        p.add_run(desc)

    doc.add_page_break()

    # ═══════════════════════════════════════════
    # 13. 별첨: 포트폴리오
    # ═══════════════════════════════════════════
    heading("13. 별첨: 포트폴리오")
    doc.add_paragraph(
        "Skyline Works는 다양한 산업 분야의 웹/앱 프로젝트를 성공적으로 수행해왔습니다. "
        "아래는 본 프로젝트와 유사한 성격의 대표 프로젝트입니다."
    )
    doc.add_paragraph()

    portfolio = [
        ("LOMA — 로케이션 마켓 플랫폼", "촬영/행사 장소 매칭 플랫폼. Next.js + Supabase 기반 풀스택 개발."),
        ("집현전 — AI 기반 교육 플랫폼", "AI 튜터링 + 콘텐츠 큐레이션 플랫폼. 지자체 대상 제안서 수주."),
        ("세바시 — 강연 플랫폼", "강연 콘텐츠 플랫폼 UI/UX 설계 및 프론트엔드 개발."),
    ]
    for title, desc in portfolio:
        p = doc.add_paragraph()
        r = p.add_run(title)
        r.bold = True
        doc.add_paragraph(desc)
        doc.add_paragraph()

    # ═══════════════════════════════════════════
    # 마무리
    # ═══════════════════════════════════════════
    doc.add_page_break()
    for _ in range(6):
        doc.add_paragraph()
    centered("감사합니다", 28, True, (0xC8, 0x96, 0x5A))
    doc.add_paragraph()
    centered("Skyline Works", 18, True, (0x1A, 0x1A, 0x1A))
    centered("스카이라인웍스", 14, False, (0x66, 0x66, 0x66))
    doc.add_paragraph()
    centered("본 제안서에 대한 문의사항이 있으시면 언제든 연락 부탁드립니다.", 11, False, (0x99, 0x99, 0x99))

    path = os.path.join(BASE_DIR, "SpaceONEZ_제안서.docx")
    doc.save(path)
    print(f"Word: {path}")


# ══════════════════════════════════════════════════════════════
#  POWERPOINT — 웹사이트 제작 제안서
# ══════════════════════════════════════════════════════════════
def generate_pptx():
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    DARK = PptRGB(0x0A, 0x0A, 0x08)
    DARK_LIGHT = PptRGB(0x17, 0x16, 0x12)
    ACCENT = PptRGB(0xC8, 0x96, 0x5A)
    WHITE = PptRGB(0xF0, 0xEB, 0xE3)
    GRAY = PptRGB(0x8A, 0x85, 0x78)
    LIGHT_GRAY = PptRGB(0x66, 0x60, 0x55)

    def dark_bg(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = DARK

    def add_textbox(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = PptPt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = alignment
        return tf

    def add_multi_text(slide, left, top, width, height, items, size=16, color=WHITE, spacing=8):
        txBox = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.font.size = PptPt(size)
            p.font.color.rgb = color
            p.space_after = PptPt(spacing)
        return tf

    def add_card(slide, x, y, w, h, title, desc, title_size=16, desc_size=13):
        shape = slide.shapes.add_shape(1, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK_LIGHT
        shape.line.fill.background()
        add_textbox(slide, x + 0.3, y + 0.3, w - 0.6, 0.5, title, title_size, ACCENT, True)
        add_textbox(slide, x + 0.3, y + 0.9, w - 0.6, h - 1.2, desc, desc_size, GRAY)

    def slide_header(slide, num, title):
        add_textbox(slide, 0.8, 0.5, 3, 0.4, f"{num:02d}", 12, ACCENT, True)
        add_textbox(slide, 0.8, 0.9, 10, 1, title, 36, WHITE, True)

    # ═══════════════ Slide 1: Cover ═══════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    add_textbox(slide, 1, 0.8, 11, 0.5, "SKYLINE WORKS", 14, ACCENT, True)
    add_textbox(slide, 1, 1.5, 11, 0.5, "웹사이트 제작 제안서", 16, GRAY)
    add_textbox(slide, 1, 2.5, 11, 2.5, "스페이스원지 &\n올드비트 영도", 52, WHITE, True)
    add_textbox(slide, 1, 5.2, 11, 0.8, "부산 영도의 복합문화공간 공식 웹사이트 구축 프로젝트", 20, GRAY)
    add_textbox(slide, 1, 6.4, 11, 0.5, "2025. 04  |  Skyline Works", 14, LIGHT_GRAY)

    # ═══════════════ Slide 2: 제안 개요 ═══════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    slide_header(slide, 1, "제안 개요")
    items = [
        "프로젝트  스페이스원지 & 올드비트 영도 공식 웹사이트 구축",
        "유형  신규 웹사이트 기획 · 디자인 · 개발",
        "위치  부산광역시 영도구 봉래나루로 214 1층",
        "공간  복합문화공간 (레스토랑 + 갤러리 + 공연장 + 워크숍)",
        "웹사이트  반응형 8개 페이지 (홈, 소개, 소식, 행사, 공간, 예약)",
        "예약  네이버 예약 외부 링크 연동",
        "기간  착수 후 5주",
    ]
    add_multi_text(slide, 0.8, 2.3, 11, 4.5, items, 18, GRAY, 12)

    # ═══════════════ Slide 3: 제안사 소개 ═══════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    slide_header(slide, 2, "제안사 소개")
    add_textbox(slide, 0.8, 2.2, 5.5, 1,
        "Skyline Works\n스카이라인웍스", 28, WHITE, True)
    add_textbox(slide, 0.8, 3.8, 5.5, 1,
        '"예쁜 웹사이트는 필요 없습니다.\n우리는 팔리는 논리를 만듭니다."', 16, GRAY)

    strengths = [
        "Solid Logic — 설득하는 비즈니스 시나리오 설계",
        "Irresistible Flow — 전환을 유도하는 UX 흐름",
        "Visual Persuasion — 논리를 각인시키는 디자인",
        "Communication Bridge — 비즈니스 ↔ 기술 통역",
        "Operational Intelligence — 런칭 이후까지 고려",
    ]
    add_multi_text(slide, 7, 2.2, 5.5, 4.5, strengths, 16, GRAY, 14)

    # ═══════════════ Slide 4: 프로젝트 이해 ═══════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    slide_header(slide, 3, "프로젝트 이해")
    add_textbox(slide, 0.8, 2.2, 5.5, 2,
        "영도의 역사적 산업 건물을 리노베이션한\n"
        "복합문화공간의 브랜드 가치를\n"
        "온라인에 효과적으로 전달",
        20, WHITE)

    challenges = [
        "두 브랜드 통합 아이덴티티 구현",
        "공간/메뉴/행사 통합 정보 허브",
        "네이버 예약 → 예약 전환 최적화",
        "온라인 대관 문의 채널 확보",
        "전시/공연/워크숍 지속 홍보",
    ]
    add_multi_text(slide, 7, 2.2, 5.5, 4, challenges, 18, GRAY, 14)

    # ═══════════════ Slide 5: 사이트 구조 ═══════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    slide_header(slide, 4, "사이트 구조")

    pages = [
        ("홈  /", "히어로 + 공간소개\n행사 + 소식"),
        ("소개  /about", "스토리 · 철학\n통계"),
        ("소식  /news", "카드 그리드\n공지/소식/미디어"),
        ("행사  /events", "진행중/예정/종료\n전시/공연/워크숍"),
        ("공간  /space", "레스토랑 · 갤러리\n루프탑 · 워크숍룸"),
        ("예약  /reservation", "네이버 예약\n대관 문의 폼"),
    ]
    for i, (title, desc) in enumerate(pages):
        col = i % 3
        row = i // 3
        x = 0.8 + col * 4.1
        y = 2.3 + row * 2.4
        add_card(slide, x, y, 3.7, 2, title, desc)

    # ═══════════════ Slide 6: 디자인 컨셉 ═══════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    slide_header(slide, 5, "디자인 컨셉")
    add_textbox(slide, 0.8, 2, 11, 0.8,
        '"Warm Industrial" — 영도 산업유산의 따뜻한 울림, 패럴랙스 스크롤, 대형 타이포', 20, GRAY)

    # Color swatches
    colors_data = [
        ("#0A0A08", "웜 블랙\n배경", DARK),
        ("#C8965A", "앤틱 브래스\nCTA/강조", ACCENT),
        ("#F0EBE3", "웜 크림\n메인 텍스트", WHITE),
        ("#8A8578", "웜 그레이\n서브 텍스트", GRAY),
    ]
    for i, (code, name, rgb) in enumerate(colors_data):
        x = 0.8 + i * 3
        shape = slide.shapes.add_shape(1, PptInches(x), PptInches(3.2), PptInches(2.6), PptInches(1.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb
        shape.line.fill.background()
        txt_color = DARK if i >= 2 else WHITE
        add_textbox(slide, x + 0.2, 3.4, 2.2, 0.4, code, 13, txt_color, True)
        add_textbox(slide, x + 0.2, 3.9, 2.2, 0.6, name, 11, txt_color)

    anims = [
        "fade-up — 아래에서 페이드인   |   reveal-scale — 스케일업",
        "reveal-left/right — 좌/우 슬라이드   |   패럴랙스 — 스크롤 기반 이동",
        "hero-zoom — 8초간 배경 줌아웃   |   Lenis — 스무스 스크롤",
    ]
    add_multi_text(slide, 0.8, 5, 11, 2, anims, 14, GRAY)

    # ═══════════════ Slide 7: 홈페이지 구성 ═══════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    slide_header(slide, 6, "홈페이지 구성")

    sections = [
        ("HERO", "120vh 패럴랙스\n줌인 애니메이션\n대형 타이포 (9xl)"),
        ("공간 소개", "풀스크린 x 2\n패럴랙스 배경\n좌/우 슬라이드인"),
        ("인터루드", "풀블리드 인용문\n스케일업 리빌\n'오래된 것의 울림'"),
        ("행사·전시", "21:9 피처드 카드\n+ 그리드 레이아웃\n상태 뱃지"),
        ("소식", "에디토리얼 리스트\n번호 + 호버 밀림\n카테고리 태그"),
    ]
    for i, (title, desc) in enumerate(sections):
        x = 0.5 + i * 2.5
        add_card(slide, x, 2.2, 2.2, 4.5, title, desc, 14, 12)

    # ═══════════════ Slide 8: 기술 스택 ═══════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    slide_header(slide, 7, "기술 스택")

    tech = [
        ("Next.js 16", "App Router\nSSG 정적 생성\nMetadata SEO"),
        ("React 19", "Server Components\nTypeScript\n컴포넌트 아키텍처"),
        ("Tailwind v4", "유틸리티 CSS\n반응형 디자인\n커스텀 테마"),
        ("Vercel 배포", "자동 빌드/배포\nCDN + SSL\n커스텀 도메인"),
    ]
    for i, (title, desc) in enumerate(tech):
        x = 0.8 + i * 3.1
        add_card(slide, x, 2.3, 2.7, 3.5, title, desc, 18, 14)

    # ═══════════════ Slide 9: 핵심 기능 ═══════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    slide_header(slide, 8, "핵심 기능")

    features = [
        ("네이버 예약 연동", "런치/디너 코스 메뉴 카드\nCTA 버튼 → 네이버 예약\n행사별 예약 유도"),
        ("대관 문의 폼", "이름, 연락처, 이메일\n희망일자, 행사유형(7종)\n접수 확인 UI"),
        ("반응형 디자인", "320px ~ 1920px 완전 대응\n모바일 햄버거 메뉴\n터치 최적화"),
        ("SEO 최적화", "Next.js Metadata API\nSSG 정적 생성\nOG 태그 지원"),
    ]
    for i, (title, desc) in enumerate(features):
        col = i % 2
        row = i // 2
        x = 0.8 + col * 6.2
        y = 2.3 + row * 2.5
        add_card(slide, x, y, 5.5, 2.1, title, desc)

    # ═══════════════ Slide 10: 추진 일정 ═══════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    slide_header(slide, 9, "추진 일정 (5주)")

    timeline = [
        ("W1", "프로젝트 셋업", "환경 구성\n디자인 토큰\nHeader/Footer"),
        ("W2", "홈페이지", "히어로, 공간소개\n인터루드\n행사, 소식"),
        ("W3", "서브 페이지", "소개, 소식, 행사\n공간, 예약\n전체 시안 리뷰"),
        ("W4", "콘텐츠 + QA", "실제 사진/텍스트\n크로스브라우저\n성능 최적화"),
        ("W5", "배포 + 런칭", "Vercel 배포\n도메인 연결\n인수인계"),
    ]
    for i, (week, title, desc) in enumerate(timeline):
        x = 0.5 + i * 2.5
        # Dot
        shape = slide.shapes.add_shape(9, PptInches(x + 0.9), PptInches(2.8), PptInches(0.3), PptInches(0.3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT
        shape.line.fill.background()
        # Line
        if i < len(timeline) - 1:
            shape = slide.shapes.add_shape(1, PptInches(x + 1.2), PptInches(2.9), PptInches(2.2), PptInches(0.05))
            shape.fill.solid()
            shape.fill.fore_color.rgb = PptRGB(0x33, 0x33, 0x33)
            shape.line.fill.background()

        add_textbox(slide, x + 0.3, 3.3, 2, 0.4, week, 14, ACCENT, True, PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.3, 3.8, 2, 0.5, title, 16, WHITE, True, PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.3, 4.4, 2, 1.5, desc, 12, GRAY, alignment=PP_ALIGN.CENTER)

    # ═══════════════ Slide 11: 견적 ═══════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    slide_header(slide, 10, "견적 및 비용")

    # Left: 항목별 비용
    cost_items = [
        "기획 및 설계                1,500,000원",
        "UI/UX 디자인              2,500,000원",
        "프론트엔드 개발           4,000,000원",
        "콘텐츠 적용                   500,000원",
        "테스트 및 배포               500,000원",
        "─────────────────────",
        "합계 (부가세 별도)       9,000,000원",
        "총 금액 (VAT 포함)      9,900,000원",
    ]
    add_multi_text(slide, 0.8, 2.3, 6, 5, cost_items, 16, GRAY, 10)

    # Right: 결제 조건
    add_textbox(slide, 7.5, 2.3, 5, 0.5, "결제 조건", 20, WHITE, True)
    payment = [
        "계약금  40%  3,960,000원  (계약 체결 시)",
        "중도금  30%  2,970,000원  (시안 확인 후)",
        "잔  금  30%  2,970,000원  (최종 납품 후)",
    ]
    add_multi_text(slide, 7.5, 3.2, 5, 3, payment, 15, GRAY, 16)

    add_textbox(slide, 7.5, 5.2, 5, 0.5, "* 상세 협의를 통해 조정 가능", 12, LIGHT_GRAY)

    # ═══════════════ Slide 12: 유지보수 ═══════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    slide_header(slide, 11, "유지보수 & 운영 지원")

    # 무상
    add_card(slide, 0.8, 2.3, 5.5, 4,
        "무상 유지보수 (3개월)",
        "납품 후 3개월간 무상 하자보수\n"
        "버그 수정, 브라우저 호환성 대응\n"
        "경미한 텍스트/이미지 수정 (월 2회)\n\n"
        "콘텐츠 업데이트 가이드 제공\n"
        "JSON 수정 매뉴얼\n"
        "이미지 규격 가이드",
        18, 14)

    # 유상
    add_card(slide, 7, 2.3, 5.5, 4,
        "유상 유지보수 (선택)",
        "Basic  200,000원/월\n"
        "  월 4회 콘텐츠 업데이트 + 모니터링\n\n"
        "Standard  400,000원/월\n"
        "  월 8회 업데이트 + 기능 개선 + 리포트\n\n"
        "Premium  별도 협의\n"
        "  상시 대응 + 신규 기능 + SEO 관리",
        18, 14)

    # ═══════════════ Slide 13: 기대 효과 ═══════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    slide_header(slide, 12, "기대 효과")

    effects = [
        "브랜드 인지도 향상      시네마틱 웹사이트로 영도 대표 문화공간 이미지 구축",
        "방문자 유입 증가         SEO + SNS 공유를 통한 온라인 노출 확대",
        "예약 전환율 향상         네이버 예약 원스톱 경험 제공",
        "대관 문의 효율화         온라인 폼을 통한 체계적 접수 및 관리",
        "문화 프로그램 홍보      전시/공연/워크숍 지속 업데이트로 재방문 유도",
        "운영 효율 극대화         JSON 기반 콘텐츠 관리로 간편 업데이트",
    ]
    add_multi_text(slide, 0.8, 2.3, 11, 4.5, effects, 18, GRAY, 14)

    # ═══════════════ Slide 14: Thank you ═══════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    add_textbox(slide, 0, 2, 13.333, 1, "감사합니다", 52, WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, 0, 3.5, 13.333, 0.8,
        "Skyline Works  |  스카이라인웍스", 22, ACCENT, False, PP_ALIGN.CENTER)
    add_textbox(slide, 0, 4.5, 13.333, 0.8,
        "스페이스원지 & 올드비트 영도 웹사이트 구축 프로젝트", 16, GRAY, False, PP_ALIGN.CENTER)
    add_textbox(slide, 0, 5.8, 13.333, 0.5,
        "본 제안서에 대한 문의사항이 있으시면 언제든 연락 부탁드립니다.", 13, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    path = os.path.join(BASE_DIR, "SpaceONEZ_제안서.pptx")
    prs.save(path)
    print(f"PPTX: {path}")


if __name__ == "__main__":
    generate_word()
    generate_pptx()
