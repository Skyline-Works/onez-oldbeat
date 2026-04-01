"""Generate Word + PowerPoint documents for Space ONE.Z & Old Beat Yeongdo project."""

from docx import Document
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt, Emu
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))


# ══════════════════════════════════════════════════════════════
#  WORD DOCUMENT
# ══════════════════════════════════════════════════════════════
def generate_word():
    doc = Document()

    style = doc.styles["Normal"]
    style.font.size = Pt(11)
    style.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
    style.paragraph_format.line_spacing = 1.5

    def heading(text, level=1):
        h = doc.add_heading(text, level=level)
        for r in h.runs:
            r.font.color.rgb = RGBColor(0x05, 0x05, 0x05)

    def centered(text, size=11, bold=False, color=None, italic=False):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = p.add_run(text)
        r.font.size = Pt(size)
        r.bold = bold
        r.italic = italic
        if color:
            r.font.color.rgb = RGBColor(*color)

    def add_table(headers, rows):
        t = doc.add_table(rows=len(rows) + 1, cols=len(headers))
        t.style = "Light Grid Accent 1"
        for i, h in enumerate(headers):
            t.rows[0].cells[i].text = h
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                t.rows[ri + 1].cells[ci].text = val
        doc.add_paragraph()

    # ── Cover ──
    for _ in range(6):
        doc.add_paragraph()
    centered("스페이스원지 & 올드비트 영도", 28, True, (0xC8, 0x96, 0x5A))
    centered("웹사이트 구축 프로젝트 기획서", 18, False, (0x33, 0x33, 0x33))
    doc.add_paragraph()
    centered("Space ONE.Z & Old Beat Yeongdo\nHomepage Development Plan", 12, False, (0x99, 0x99, 0x99), True)
    for _ in range(4):
        doc.add_paragraph()
    for label, value in [
        ("프로젝트명", "ONE.Z & Old Beat Yeongdo 웹사이트 구축"),
        ("작성일", "2025년 4월"),
        ("버전", "v1.0"),
    ]:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = p.add_run(f"{label}: ")
        r.font.size = Pt(10)
        r.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        r = p.add_run(value)
        r.font.size = Pt(10)
        r.bold = True
    doc.add_page_break()

    # ── TOC ──
    heading("목차")
    for item in [
        "1. 프로젝트 개요", "2. 프로젝트 배경 및 목적", "3. 대상 및 타겟",
        "4. 사이트 구조 (사이트맵)", "5. 페이지별 상세 기획", "6. 디자인 컨셉",
        "7. 기술 스택", "8. 핵심 기능 명세", "9. 콘텐츠 관리 방안",
        "10. 일정 계획", "11. 기대 효과",
    ]:
        p = doc.add_paragraph(item)
        p.paragraph_format.space_after = Pt(4)
    doc.add_page_break()

    # ── 1. 프로젝트 개요 ──
    heading("1. 프로젝트 개요")
    add_table(["항목", "내용"], [
        ["프로젝트명", "Space ONE.Z & Old Beat Yeongdo 웹사이트 구축"],
        ["프로젝트 유형", "신규 웹사이트 개발"],
        ["위치", "부산광역시 영도구 봉래나루로 214 1층"],
        ["공간 유형", "복합문화공간 (레스토랑 + 갤러리 + 공연장 + 워크숍)"],
        ["브랜드", "스페이스원지 (Space ONE.Z) / 올드비트 영도 (Old Beat Yeongdo)"],
        ["웹사이트 구조", "하나의 통합 사이트 (두 브랜드를 섹션으로 구분)"],
        ["예약 연동", "네이버 예약 외부 링크"],
    ])

    # ── 2. 배경 및 목적 ──
    heading("2. 프로젝트 배경 및 목적")
    heading("2.1 배경", 2)
    doc.add_paragraph(
        "부산 영도는 대한민국 근현대 조선산업의 중심지로, 최근 도시재생 사업과 함께 "
        "문화예술 공간으로 변모하고 있습니다. 스페이스원지와 올드비트 영도는 영도의 "
        "역사적 산업 건물을 리노베이션하여 레스토랑, 갤러리, 공연장, 워크숍 공간을 "
        "갖춘 복합문화공간으로 탄생했습니다."
    )
    doc.add_paragraph(
        "이 공간의 가치와 프로그램을 효과적으로 전달하고, 온라인을 통한 예약과 "
        "대관 문의를 가능하게 하는 공식 웹사이트 구축이 필요합니다."
    )
    heading("2.2 목적", 2)
    for p in [
        "브랜드 아이덴티티 확립: 스페이스원지와 올드비트 영도의 독창적인 브랜드 이미지를 온라인에 구현",
        "정보 허브: 공간 소개, 메뉴, 행사/전시 일정, 소식 등 통합 정보 제공",
        "예약 채널: 네이버 예약 연동을 통한 원활한 레스토랑/행사 예약 유도",
        "대관 문의: 공간 대관 희망자를 위한 온라인 문의 채널 확보",
        "문화 플랫폼: 진행 중인 전시, 공연, 워크숍 등 문화 콘텐츠의 지속적 홍보",
    ]:
        doc.add_paragraph(p, style="List Bullet")
    doc.add_page_break()

    # ── 3. 타겟 ──
    heading("3. 대상 및 타겟")
    add_table(["타겟 그룹", "설명"], [
        ["미식 애호가", "영도 로컬 식재료 기반 코스 다이닝에 관심 있는 2030~4050세대"],
        ["문화예술 관심층", "전시, 공연, 워크숍에 관심 있는 예술 애호가 및 크리에이터"],
        ["관광객", "부산 여행 중 영도의 문화 명소를 찾는 국내외 관광객"],
        ["기업/단체", "기업 행사, 파티, 촬영 등 공간 대관을 원하는 기업 및 단체"],
        ["지역 커뮤니티", "영도 및 부산 지역 주민, 로컬 아티스트"],
    ])

    # ── 4. 사이트맵 ──
    heading("4. 사이트 구조 (사이트맵)")
    add_table(["경로", "페이지명", "주요 콘텐츠"], [
        ["/", "홈", "히어로 + 공간소개 + 인터루드 + 행사 + 소식"],
        ["/about", "소개", "스토리, 철학, 통계"],
        ["/news", "소식", "소식 카드 그리드"],
        ["/news/[slug]", "소식 상세", "개별 소식 본문"],
        ["/events", "행사/전시", "진행중/예정/종료 분류 목록"],
        ["/events/[slug]", "행사 상세", "행사 본문 + 네이버 예약 CTA"],
        ["/space", "공간 안내", "레스토랑, 갤러리, 루프탑, 워크숍룸"],
        ["/reservation", "예약/대관", "네이버 예약 + 대관 문의 폼"],
    ])
    doc.add_page_break()

    # ── 5. 페이지별 상세 ──
    heading("5. 페이지별 상세 기획")
    pages_detail = {
        "5.1 홈페이지": [
            "풀스크린 히어로: 120vh, 패럴랙스 배경, 줌인 애니메이션, 최대 9xl 타이포",
            "공간 소개: 풀스크린 시네마틱 섹션 x2, 패럴랙스 배경 + 좌/우 슬라이드인",
            "인터루드: 풀블리드 인용문, 스케일업 리빌",
            "행사: 21:9 피처드 카드 + 그리드",
            "소식: 에디토리얼 리스트 (번호 + 호버 밀림)",
        ],
        "5.2 소개": ["영도 스토리", "운영 철학", "통계 섹션", "CTA 버튼"],
        "5.3 소식": ["3열 카드 그리드", "카테고리: 공지/소식/미디어", "상세 페이지"],
        "5.4 행사": ["상태별 분류 (진행중/예정/종료)", "카테고리: 전시/공연/워크숍/팝업", "네이버 예약 CTA"],
        "5.5 공간 안내": ["4개 공간 소개", "이미지 그리드 + 특징 태그", "좌우 교차 레이아웃"],
        "5.6 예약/대관": ["네이버 예약 링크 + 메뉴 카드", "대관 문의 폼 (7개 필드)", "접수 확인 UI"],
    }
    for title, items in pages_detail.items():
        heading(title, 2)
        for item in items:
            doc.add_paragraph(item, style="List Bullet")
    doc.add_page_break()

    # ── 6. 디자인 ──
    heading("6. 디자인 컨셉")
    heading("6.1 디자인 방향", 2)
    doc.add_paragraph(
        "Warm Industrial 시네마틱 디자인. 풀블리드 이미지, 패럴랙스, "
        "대형 타이포, 앤틱 브래스 톤으로 영도 산업유산의 따뜻한 감성 표현."
    )
    heading("6.2 컬러 팔레트", 2)
    add_table(["요소", "색상 코드", "설명"], [
        ["배경", "#0A0A08", "웜 블랙"],
        ["액센트", "#C8965A", "앤틱 브래스"],
        ["텍스트", "#F0EBE3", "웜 크림"],
        ["서브 텍스트", "#8A8578", "웜 그레이"],
    ])
    heading("6.3 애니메이션", 2)
    add_table(["이름", "동작", "사용 위치"], [
        ["reveal", "아래에서 페이드인", "기본 섹션"],
        ["reveal-scale", "스케일업 페이드인", "인터루드, 피처드"],
        ["reveal-left/right", "좌/우 슬라이드인", "공간 소개"],
        ["패럴랙스", "스크롤 기반 translateY", "히어로, 공간 배경"],
    ])

    # ── 7. 기술 스택 ──
    heading("7. 기술 스택")
    add_table(["영역", "기술"], [
        ["프레임워크", "Next.js 16 (App Router)"],
        ["UI", "React 19 + TypeScript"],
        ["스타일링", "Tailwind CSS v4"],
        ["애니메이션", "CSS Transitions + IntersectionObserver"],
        ["콘텐츠", "JSON 파일 기반"],
        ["예약", "네이버 예약 외부 링크"],
        ["배포", "Vercel"],
    ])

    # ── 8. 기능 명세 ──
    heading("8. 핵심 기능 명세")
    heading("8.1 네이버 예약", 2)
    doc.add_paragraph("레스토랑/행사 예약은 네이버 예약 외부 링크로 연동. 메뉴 카드 + CTA 버튼 제공.")
    heading("8.2 대관 문의 폼", 2)
    for f in ["이름", "연락처", "이메일", "희망 일자", "행사 유형", "예상 인원", "상세 내용"]:
        doc.add_paragraph(f, style="List Bullet")
    heading("8.3 반응형", 2)
    doc.add_paragraph("320px~1920px 완전 반응형. 모바일 햄버거 메뉴.")
    heading("8.4 SEO", 2)
    doc.add_paragraph("Next.js Metadata API + SSG로 검색엔진 최적화.")
    doc.add_page_break()

    # ── 9. 콘텐츠 관리 ──
    heading("9. 콘텐츠 관리 방안")
    doc.add_paragraph("JSON 파일 기반. Git push 시 자동 반영.")
    add_table(["데이터", "주요 필드"], [
        ["news.json", "slug, title, date, thumbnail, summary, content, category"],
        ["events.json", "slug, title, startDate, endDate, status, thumbnail, description, location, category"],
        ["spaces.json", "id, name, description, images, capacity, features"],
    ])

    # ── 10. 일정 ──
    heading("10. 일정 계획")
    add_table(["기간", "작업", "상세"], [
        ["1주차", "프로젝트 셋업", "환경 구성, 디자인 토큰, Header/Footer"],
        ["2주차", "홈페이지 구현", "히어로, 공간소개, 행사, 소식"],
        ["3주차", "서브 페이지", "소개, 소식, 행사, 공간, 예약"],
        ["4주차", "콘텐츠 + QA", "실제 사진/텍스트, 크로스브라우저"],
        ["5주차", "배포 + 런칭", "Vercel, 도메인, 최종 점검"],
    ])

    # ── 11. 기대 효과 ──
    heading("11. 기대 효과")
    for e in [
        "브랜드 인지도 향상: 시네마틱 웹사이트로 영도 대표 문화공간 이미지 구축",
        "방문자 유입 증가: SEO + SNS 공유를 통한 온라인 노출 확대",
        "예약 전환율 향상: 네이버 예약 원스톱 경험",
        "대관 문의 효율화: 온라인 폼 체계적 접수",
        "문화 프로그램 홍보: 지속적 업데이트로 재방문 유도",
        "운영 효율: JSON 기반 간편 업데이트",
    ]:
        doc.add_paragraph(e, style="List Bullet")

    doc.add_paragraph()
    centered("- End of Document -", 10, False, (0x99, 0x99, 0x99), True)

    path = os.path.join(BASE_DIR, "SpaceONEZ_기획서.docx")
    doc.save(path)
    print(f"Word: {path}")


# ══════════════════════════════════════════════════════════════
#  POWERPOINT
# ══════════════════════════════════════════════════════════════
def generate_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    DARK = PptRGB(0x0A, 0x0A, 0x08)
    DARK_LIGHT = PptRGB(0x17, 0x16, 0x12)
    ACCENT = PptRGB(0xC8, 0x96, 0x5A)
    WHITE = PptRGB(0xF0, 0xEB, 0xE3)
    GRAY = PptRGB(0x8A, 0x85, 0x78)

    def dark_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = DARK

    def add_textbox(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = PptPt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = alignment
        return tf

    def add_multi_text(slide, left, top, width, height, items, size=16, color=WHITE):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = PptPt(size)
            p.font.color.rgb = color
            p.space_after = PptPt(8)
        return tf

    # ── Slide 1: Cover ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    dark_bg(slide)
    add_textbox(slide, 1, 1.5, 11, 0.5, "SPACE ONE.Z & OLD BEAT YEONGDO", 14, ACCENT, True)
    add_textbox(slide, 1, 2.3, 11, 2, "스페이스원지 &\n올드비트 영도", 48, WHITE, True)
    add_textbox(slide, 1, 4.8, 11, 1, "웹사이트 구축 프로젝트 기획서", 24, GRAY)
    add_textbox(slide, 1, 6.2, 11, 0.5, "2025. 04  |  v1.0", 14, GRAY)

    # ── Slide 2: 프로젝트 개요 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    add_textbox(slide, 0.8, 0.5, 3, 0.4, "01", 12, ACCENT, True)
    add_textbox(slide, 0.8, 0.9, 6, 1, "프로젝트 개요", 36, WHITE, True)
    items = [
        "프로젝트명  Space ONE.Z & Old Beat Yeongdo 웹사이트",
        "위치  부산광역시 영도구 봉래나루로 214 1층",
        "공간 유형  복합문화공간 (레스토랑 + 갤러리 + 공연장 + 워크숍)",
        "웹사이트  하나의 통합 사이트",
        "예약 연동  네이버 예약 외부 링크",
        "기술 스택  Next.js 16 + React 19 + Tailwind v4",
    ]
    add_multi_text(slide, 0.8, 2.2, 11, 4.5, items, 18, GRAY)

    # ── Slide 3: 배경 & 목적 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    add_textbox(slide, 0.8, 0.5, 3, 0.4, "02", 12, ACCENT, True)
    add_textbox(slide, 0.8, 0.9, 6, 1, "배경 & 목적", 36, WHITE, True)
    add_textbox(slide, 0.8, 2.2, 5.5, 2.5,
        "부산 영도의 역사적 산업 건물을 리노베이션한\n"
        "복합문화공간의 가치와 프로그램을\n"
        "효과적으로 전달하는 공식 웹사이트 구축",
        20, WHITE)
    goals = [
        "브랜드 아이덴티티 온라인 구현",
        "공간/메뉴/행사 통합 정보 허브",
        "네이버 예약 연동 → 예약 전환",
        "온라인 대관 문의 채널 확보",
        "전시/공연/워크숍 지속 홍보",
    ]
    add_multi_text(slide, 7, 2.2, 5.5, 4, goals, 18, GRAY)

    # ── Slide 4: 사이트맵 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    add_textbox(slide, 0.8, 0.5, 3, 0.4, "03", 12, ACCENT, True)
    add_textbox(slide, 0.8, 0.9, 6, 1, "사이트 구조", 36, WHITE, True)

    pages = [
        ("홈 /", "히어로 + 공간소개\n+ 행사 + 소식"),
        ("소개 /about", "스토리 · 철학\n통계"),
        ("소식 /news", "카드 그리드\n공지/소식/미디어"),
        ("행사 /events", "진행중/예정/종료\n전시/공연/워크숍"),
        ("공간 /space", "레스토랑 · 갤러리\n루프탑 · 워크숍룸"),
        ("예약 /reservation", "네이버 예약\n대관 문의 폼"),
    ]
    for i, (title, desc) in enumerate(pages):
        col = i % 3
        row = i // 3
        x = 0.8 + col * 4.1
        y = 2.3 + row * 2.4

        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(3.7), Inches(2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK_LIGHT
        shape.line.fill.background()

        add_textbox(slide, x + 0.3, y + 0.3, 3.2, 0.5, title, 16, ACCENT, True)
        add_textbox(slide, x + 0.3, y + 0.9, 3.2, 1, desc, 13, GRAY)

    # ── Slide 5: 디자인 컨셉 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    add_textbox(slide, 0.8, 0.5, 3, 0.4, "04", 12, ACCENT, True)
    add_textbox(slide, 0.8, 0.9, 8, 1, "디자인 컨셉", 36, WHITE, True)
    add_textbox(slide, 0.8, 2, 11, 0.8,
        '"Warm Industrial" — 영도 산업유산의 따뜻한 감성, 몰입감 있는 경험', 20, GRAY)

    # Color swatches
    colors_data = [
        ("#0A0A08", "웜 블랙", DARK),
        ("#C8965A", "앤틱 브래스", ACCENT),
        ("#F0EBE3", "웜 크림", WHITE),
        ("#8A8578", "웜 그레이", GRAY),
    ]
    for i, (code, name, rgb) in enumerate(colors_data):
        x = 0.8 + i * 3
        shape = slide.shapes.add_shape(1, Inches(x), Inches(3.2), Inches(2.6), Inches(1.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb
        shape.line.fill.background()
        txt_color = DARK if i >= 2 else WHITE
        add_textbox(slide, x + 0.2, Inches(3.4).inches, 2.2, 0.4, code, 12, txt_color, True)
        add_textbox(slide, x + 0.2, Inches(3.8).inches, 2.2, 0.4, name, 11, txt_color)

    # Animations
    anims = [
        "reveal — 아래에서 페이드인",
        "reveal-scale — 스케일업 페이드인",
        "reveal-left/right — 좌/우 슬라이드인",
        "패럴랙스 — 스크롤 기반 배경 이동",
        "히어로 줌인 — 8초간 1.15→1 스케일",
    ]
    add_multi_text(slide, 0.8, 4.8, 11, 2, anims, 15, GRAY)

    # ── Slide 6: 홈페이지 구성 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    add_textbox(slide, 0.8, 0.5, 3, 0.4, "05", 12, ACCENT, True)
    add_textbox(slide, 0.8, 0.9, 8, 1, "홈페이지 구성", 36, WHITE, True)

    sections = [
        ("HERO", "120vh 패럴랙스\n줌인 애니메이션\n대형 타이포 (9xl)"),
        ("공간 소개", "풀스크린 x 2\n패럴랙스 배경\n좌/우 슬라이드인"),
        ("인터루드", "풀블리드 인용문\n스케일업 리빌\n'오래된 것의 새로운 울림'"),
        ("행사·전시", "21:9 피처드 카드\n+ 2열 그리드\n상태 뱃지"),
        ("소식", "에디토리얼 리스트\n번호 + 호버 밀림\n카테고리 태그"),
    ]
    for i, (title, desc) in enumerate(sections):
        x = 0.5 + i * 2.5
        shape = slide.shapes.add_shape(1, Inches(x), Inches(2.2), Inches(2.2), Inches(4.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK_LIGHT
        shape.line.fill.background()
        add_textbox(slide, x + 0.2, 2.5, 1.8, 0.5, title, 14, ACCENT, True)
        add_textbox(slide, x + 0.2, 3.2, 1.8, 3, desc, 12, GRAY)

    # ── Slide 7: 기술 스택 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    add_textbox(slide, 0.8, 0.5, 3, 0.4, "06", 12, ACCENT, True)
    add_textbox(slide, 0.8, 0.9, 6, 1, "기술 스택", 36, WHITE, True)

    tech = [
        ("Next.js 16", "App Router\nSSG 정적 생성\nMetadata SEO"),
        ("React 19", "Server Components\nTypeScript\n컴포넌트 아키텍처"),
        ("Tailwind v4", "유틸리티 CSS\n반응형 디자인\n커스텀 테마"),
        ("CSS Animations", "GPU 가속 transition\nIntersectionObserver\n패럴랙스 스크롤"),
    ]
    for i, (title, desc) in enumerate(tech):
        x = 0.8 + i * 3.1
        shape = slide.shapes.add_shape(1, Inches(x), Inches(2.3), Inches(2.7), Inches(3.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK_LIGHT
        shape.line.fill.background()
        add_textbox(slide, x + 0.3, 2.6, 2.2, 0.5, title, 18, ACCENT, True)
        add_textbox(slide, x + 0.3, 3.3, 2.2, 2, desc, 14, GRAY)

    # ── Slide 8: 핵심 기능 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    add_textbox(slide, 0.8, 0.5, 3, 0.4, "07", 12, ACCENT, True)
    add_textbox(slide, 0.8, 0.9, 6, 1, "핵심 기능", 36, WHITE, True)

    features = [
        ("네이버 예약 연동", "런치/디너 코스 메뉴 카드\nCTA 버튼 → 네이버 예약\n행사 상세에서도 예약 유도"),
        ("대관 문의 폼", "이름, 연락처, 이메일\n희망일자, 행사유형, 인원\n접수 확인 UI"),
        ("반응형 디자인", "320px ~ 1920px\n모바일 햄버거 메뉴\n1열 레이아웃 전환"),
        ("SEO 최적화", "Next.js Metadata API\nSSG 정적 생성\n페이지별 title/description"),
    ]
    for i, (title, desc) in enumerate(features):
        col = i % 2
        row = i // 2
        x = 0.8 + col * 6.2
        y = 2.3 + row * 2.5
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(5.5), Inches(2.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK_LIGHT
        shape.line.fill.background()
        add_textbox(slide, x + 0.3, y + 0.3, 5, 0.4, title, 16, ACCENT, True)
        add_textbox(slide, x + 0.3, y + 0.8, 5, 1.2, desc, 13, GRAY)

    # ── Slide 9: 일정 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    add_textbox(slide, 0.8, 0.5, 3, 0.4, "08", 12, ACCENT, True)
    add_textbox(slide, 0.8, 0.9, 6, 1, "일정 계획", 36, WHITE, True)

    timeline = [
        ("W1", "프로젝트 셋업", "환경, 디자인 토큰\nHeader/Footer"),
        ("W2", "홈페이지", "히어로, 공간소개\n행사, 소식"),
        ("W3", "서브 페이지", "소개, 소식, 행사\n공간, 예약"),
        ("W4", "콘텐츠 + QA", "실제 사진/텍스트\n크로스브라우저"),
        ("W5", "배포 + 런칭", "Vercel 배포\n도메인, 최종 점검"),
    ]
    for i, (week, title, desc) in enumerate(timeline):
        x = 0.5 + i * 2.5
        # Timeline dot
        shape = slide.shapes.add_shape(9, Inches(x + 0.9), Inches(2.8), Inches(0.3), Inches(0.3))  # oval
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT
        shape.line.fill.background()
        # Line
        if i < len(timeline) - 1:
            shape = slide.shapes.add_shape(1, Inches(x + 1.2), Inches(2.9), Inches(2.2), Inches(0.05))
            shape.fill.solid()
            shape.fill.fore_color.rgb = PptRGB(0x33, 0x33, 0x33)
            shape.line.fill.background()

        add_textbox(slide, x + 0.3, 3.3, 2, 0.4, week, 14, ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.3, 3.8, 2, 0.5, title, 16, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.3, 4.4, 2, 1.2, desc, 12, GRAY, alignment=PP_ALIGN.CENTER)

    # ── Slide 10: 기대 효과 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    add_textbox(slide, 0.8, 0.5, 3, 0.4, "09", 12, ACCENT, True)
    add_textbox(slide, 0.8, 0.9, 6, 1, "기대 효과", 36, WHITE, True)

    effects = [
        "브랜드 인지도 향상  시네마틱 웹사이트로 영도 대표 문화공간 이미지 구축",
        "방문자 유입 증가  SEO + SNS 공유를 통한 온라인 노출 확대",
        "예약 전환율 향상  네이버 예약 원스톱 경험 제공",
        "대관 문의 효율화  온라인 폼을 통한 체계적 접수 및 관리",
        "문화 프로그램 홍보  전시/공연/워크숍 지속 업데이트로 재방문 유도",
        "운영 효율  JSON 기반 콘텐츠 관리로 간편 업데이트",
    ]
    add_multi_text(slide, 0.8, 2.3, 11, 4.5, effects, 18, GRAY)

    # ── Slide 11: Thank you ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    add_textbox(slide, 0, 2, 13.333, 1, "감사합니다", 48, WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, 0, 3.5, 13.333, 0.8,
        "Space ONE.Z & Old Beat Yeongdo", 20, ACCENT, False, PP_ALIGN.CENTER)
    add_textbox(slide, 0, 4.5, 13.333, 0.5,
        "부산광역시 영도구 봉래나루로 214 1층", 14, GRAY, False, PP_ALIGN.CENTER)

    path = os.path.join(BASE_DIR, "SpaceONEZ_기획서.pptx")
    prs.save(path)
    print(f"PPTX: {path}")


if __name__ == "__main__":
    generate_word()
    generate_pptx()
